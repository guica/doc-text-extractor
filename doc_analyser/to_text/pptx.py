from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def find_text_in_list_of_shape(shapes, df=None):
    if df is None:
        df = ''

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            find_text_in_list_of_shape(shape.shapes, df=df)

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    df = df + ' ' + run.text
                    
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    df = df + ' ' + cell.text
    return df

def convert_pptx_to_txt(path):
    prs = Presentation(path)
    result = {}
    result['full'] = ''
    for i, slide in enumerate(prs.slides):
        text_data = find_text_in_list_of_shape(slide.shapes)
        result[i+1] = text_data
        result['full'] += text_data
    
    return result